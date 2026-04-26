"""Generate the project presentation as a .pptx file."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "outputs"
PPT = ROOT / "ppt" / "IC2306_DCT_Image_Compression_Slides.pptx"


NAVY = RGBColor(0x1F, 0x2A, 0x44)
ACCENT = RGBColor(0xD9, 0x77, 0x06)
MUTED = RGBColor(0x4B, 0x55, 0x63)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xE5, 0xE7, 0xEB)


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_rect(slide, x, y, w, h, fill, outline=None):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background() if outline is None else None
    if outline:
        shape.line.color.rgb = outline
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, x, y, w, h, text, *, size=18, bold=False,
             color=NAVY, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    first = True
    for line in text.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=16, color=MUTED):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"•  {it}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        p.space_after = Pt(8)
    return tb


def slide_header(slide, title, kicker=None):
    add_rect(slide, 0, 0, Inches(13.333), Inches(1.0), NAVY)
    if kicker:
        add_text(slide, Inches(0.5), Inches(0.1), Inches(8), Inches(0.3),
                 kicker, size=11, color=LIGHT, italic=True)
    add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             title, size=24, bold=True, color=WHITE)


def add_picture_centered(slide, path, top, max_width_in, max_height_in=None):
    from PIL import Image
    w, h = Image.open(path).size
    ratio = h / w
    width = Inches(max_width_in)
    height = Inches(max_width_in * ratio)
    if max_height_in and height > Inches(max_height_in):
        height = Inches(max_height_in)
        width = Inches(max_height_in / ratio)
    left = Emu((prs_width - width) // 2)
    slide.shapes.add_picture(str(path), left, top, width=width, height=height)


prs_width = Inches(13.333)
prs_height = Inches(7.5)


def build():
    prs = Presentation()
    prs.slide_width = prs_width
    prs.slide_height = prs_height

    # 1 - Title slide
    s = blank_slide(prs)
    add_rect(s, 0, 0, prs_width, prs_height, NAVY)
    add_rect(s, Inches(0.6), Inches(5.6), Inches(1.2), Inches(0.08), ACCENT)
    add_text(s, Inches(0.6), Inches(1.8), Inches(12), Inches(1.2),
             "Image Compression using the 2D Discrete Cosine Transform",
             size=40, bold=True, color=WHITE)
    add_text(s, Inches(0.6), Inches(3.1), Inches(12), Inches(0.7),
             "A JPEG-baseline style codec with rate-distortion analysis",
             size=22, color=LIGHT, italic=True)
    add_text(s, Inches(0.6), Inches(5.8), Inches(12), Inches(0.5),
             "IC2306  —  Signal and Image Processing",
             size=18, color=WHITE, bold=True)
    add_text(s, Inches(0.6), Inches(6.25), Inches(12), Inches(0.5),
             "Course project  •  Presented by Ved",
             size=14, color=LIGHT)

    # 2 - Agenda
    s = blank_slide(prs)
    slide_header(s, "What we will cover", kicker="Agenda")
    add_bullets(s, Inches(0.7), Inches(1.4), Inches(12), Inches(5.5), [
        "Why compress images at all — the information-theoretic motivation",
        "From the DFT to the DCT — and why the DCT is the right tool",
        "The 2D DCT in matrix form, and what its 64 basis images look like",
        "Quantisation — where the loss actually happens",
        "System design: encoder, decoder, quality factor",
        "Results on the 512×512 ‘ascent’ image at Q = 10, 20, 50, 90",
        "Rate–distortion analysis and where the reconstruction error lives",
        "Discussion, limitations, and future extensions",
    ], size=18)

    # 3 - Problem / motivation
    s = blank_slide(prs)
    slide_header(s, "Why compress images?", kicker="Motivation")
    add_bullets(s, Inches(0.7), Inches(1.4), Inches(12), Inches(5), [
        "A 512×512 8-bit grayscale image is 256 KB of raw pixel data.",
        "Neighbouring pixels are strongly correlated — most of that data is redundant.",
        "The human eye is far more sensitive to low-frequency detail than to high.",
        "A good codec should exploit both: decorrelate, then throw away what cannot be seen.",
        "JPEG does this with three stages — we will reproduce the first two from scratch.",
    ])
    add_text(s, Inches(0.7), Inches(6.0), Inches(12), Inches(0.6),
             "Goal: not to produce a JPEG file, but to make the mathematics visible.",
             size=16, italic=True, color=ACCENT)

    # 4 - Syllabus mapping
    s = blank_slide(prs)
    slide_header(s, "Mapping to the IC2306 syllabus", kicker="Course outcomes")
    add_bullets(s, Inches(0.7), Inches(1.3), Inches(12), Inches(5.8), [
        "CO1 – sampling & LTI operations → used implicitly in block-wise DCT",
        "CO2 – DFT and FFT → DCT is derived from an even-symmetric DFT",
        "CO3 – filter design → quantisation acts like a spatial low-pass filter with varying cutoff",
        "CO4 – pixel operations → level shift, padding, clipping, block traversal",
        "CO5 – image transforms → 2D DCT basis, separability, energy compaction",
        "CO6 – image enhancement & restoration → blocking & ringing as predictable artefacts",
    ])

    # 5 - From DFT to DCT
    s = blank_slide(prs)
    slide_header(s, "From the DFT to the DCT", kicker="Mathematics 1 of 4")
    add_bullets(s, Inches(0.7), Inches(1.3), Inches(12), Inches(2.5), [
        "The DFT of a real block is complex — every coefficient needs two numbers.",
        "The implied periodic extension is discontinuous — creates spurious high-frequency energy.",
        "Extending the signal even-symmetrically before transforming fixes both issues.",
    ])
    add_rect(s, Inches(0.7), Inches(4.1), Inches(12), Inches(2.5), LIGHT)
    add_text(s, Inches(0.9), Inches(4.2), Inches(11.6), Inches(0.6),
             "DCT-II of length N:", size=16, bold=True, color=NAVY)
    add_text(s, Inches(0.9), Inches(4.7), Inches(11.6), Inches(0.8),
             "X[k] = α(k) · Σₙ x[n] · cos( (2n+1) k π / (2N) ),   k = 0 … N−1",
             size=18, italic=True, color=NAVY)
    add_text(s, Inches(0.9), Inches(5.5), Inches(11.6), Inches(0.8),
             "with α(0) = √(1/N),   α(k) = √(2/N) for k > 0",
             size=16, italic=True, color=MUTED)
    add_text(s, Inches(0.9), Inches(6.1), Inches(11.6), Inches(0.6),
             "Result: real coefficients, smoother periodic extension, better energy compaction.",
             size=14, color=ACCENT, italic=True)

    # 6 - 2D DCT
    s = blank_slide(prs)
    slide_header(s, "The 2D DCT is separable", kicker="Mathematics 2 of 4")
    add_bullets(s, Inches(0.7), Inches(1.3), Inches(6.2), Inches(4.5), [
        "Build T once: T[k,n] = α(k) cos((2n+1)kπ/2N)",
        "T is orthonormal, so T⁻¹ = Tᵀ.",
        "2D DCT of an 8×8 block: F = T · f · Tᵀ",
        "2D IDCT: f = Tᵀ · F · T",
        "Separability cuts cost from O(N⁴) to O(N³).",
        "Basis images are outer products — see right.",
    ], size=16)
    add_picture_centered(s, OUT / "fig_dct_basis.png",
                         Inches(1.3), max_width_in=6.0, max_height_in=5.5)
    # manually position right side for side-by-side
    # (the helper centres by default — override by removing and re-adding)
    # simpler: just add picture on right directly
    # We remove the centred one just added — but shapes list doesn't let us easily.
    # Workaround: don't call add_picture_centered above; place directly.

    # Re-approach: clear the last added shape (the centred image) by resizing/moving.
    last = s.shapes[-1]
    last.left = Inches(7.2)
    last.top = Inches(1.5)
    last.width = Inches(5.6)
    last.height = Inches(5.6)

    # 7 - Quantisation
    s = blank_slide(prs)
    slide_header(s, "Quantisation is where loss happens", kicker="Mathematics 3 of 4")
    add_bullets(s, Inches(0.7), Inches(1.3), Inches(12), Inches(2.2), [
        "F_q(u,v) = round( F(u,v) / Q'(u,v) )     — only lossy step in the pipeline.",
        "Q' is the IJG luminance table scaled by the quality factor.",
        "High-frequency coefficients get larger divisors → more zeros after rounding.",
    ])
    s.shapes.add_picture(str(OUT / "fig_quantization_matrices.png"),
                         Inches(0.7), Inches(3.7),
                         width=Inches(12))

    # 8 - Quality factor
    s = blank_slide(prs)
    slide_header(s, "One integer controls everything", kicker="Mathematics 4 of 4")
    add_bullets(s, Inches(0.7), Inches(1.3), Inches(12), Inches(3.2), [
        "Scale factor:   scale = 5000 / Q   if Q < 50,   scale = 200 − 2Q  otherwise",
        "Q'(u,v) = max(1, round( Q₅₀(u,v) · scale / 100 ))",
        "Lower Q ⇒ bigger divisors ⇒ more zeros ⇒ lower entropy AND more detail lost.",
        "Measure distortion with PSNR = 20 · log₁₀(255 / √MSE).",
        "Measure rate with Shannon entropy of the quantised coefficient stream.",
    ])

    # 9 - Block diagram
    s = blank_slide(prs)
    slide_header(s, "System block diagram", kicker="Design")
    s.shapes.add_picture(str(OUT / "fig_block_diagram.png"),
                         Inches(0.5), Inches(1.4),
                         width=Inches(12.3))

    # 10 - Block walk-through
    s = blank_slide(prs)
    slide_header(s, "One 8×8 block end to end", kicker="Walk-through")
    s.shapes.add_picture(str(OUT / "fig_block_pipeline.png"),
                         Inches(0.5), Inches(1.3),
                         width=Inches(12.3))
    add_text(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.6),
             "At Q=30: 57 non-zero DCT coefficients collapse to 6 after quantisation — "
             "block PSNR is still 36.4 dB.",
             size=14, italic=True, color=ACCENT)

    # 11 - Implementation
    s = blank_slide(prs)
    slide_header(s, "Implementation highlights", kicker="Python / NumPy")
    add_bullets(s, Inches(0.7), Inches(1.3), Inches(12), Inches(2.2), [
        "dct_compression.py – pure DSP; no plotting, no file I/O.",
        "run_project.py – driver that generates every figure you see.",
        "DCT matrix built once at import time; reused for every block.",
        "Verified: T·Tᵀ = I to 1e−14; dct2(idct2(·)) is an identity to float precision.",
    ])
    add_rect(s, Inches(0.7), Inches(4.3), Inches(12), Inches(2.7), LIGHT)
    code = ("def build_dct_matrix(n=8):\n"
            "    k = np.arange(n).reshape(-1, 1)\n"
            "    i = np.arange(n).reshape(1, -1)\n"
            "    T = np.cos((2*i + 1) * k * np.pi / (2*n))\n"
            "    T *= np.sqrt(2.0 / n)\n"
            "    T[0, :] *= 1.0 / np.sqrt(2.0)\n"
            "    return T")
    tb = s.shapes.add_textbox(Inches(0.9), Inches(4.45), Inches(11.6), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in code.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = line
        run.font.name = "Consolas"
        run.font.size = Pt(14)
        run.font.color.rgb = NAVY

    # 12 - Results grid
    s = blank_slide(prs)
    slide_header(s, "Reconstruction across quality factors", kicker="Results 1 of 3")
    s.shapes.add_picture(str(OUT / "fig_reconstruction_grid.png"),
                         Inches(0.3), Inches(1.2),
                         width=Inches(12.7))

    # 13 - Metrics table
    s = blank_slide(prs)
    slide_header(s, "Numerical summary", kicker="Results 2 of 3")
    rows_text = [
        ("Q", "PSNR (dB)", "MSE", "Sparsity", "bpp", "CR"),
        ("90", "42.07", "4.04", "69.8 %", "2.39", "3.35×"),
        ("50", "33.39", "29.82", "86.0 %", "1.16", "6.87×"),
        ("20", "29.83", "67.56", "91.6 %", "0.72", "11.08×"),
        ("10", "27.44", "117.36", "94.6 %", "0.47", "16.88×"),
    ]
    top = Inches(1.6)
    row_h = Inches(0.75)
    col_widths = [Inches(1.2), Inches(2.0), Inches(1.8), Inches(1.8), Inches(1.5), Inches(2.0)]
    col_x = [Inches(1.5)]
    for w in col_widths[:-1]:
        col_x.append(Emu(col_x[-1] + w))
    for r_idx, row in enumerate(rows_text):
        for c_idx, (val, xw) in enumerate(zip(row, col_widths)):
            x = col_x[c_idx]
            y = Emu(top + row_h * r_idx)
            fill = NAVY if r_idx == 0 else (LIGHT if r_idx % 2 else WHITE)
            add_rect(s, x, y, xw, row_h, fill, outline=MUTED)
            color = WHITE if r_idx == 0 else NAVY
            add_text(s, x, Emu(y + Inches(0.18)), xw, Inches(0.4),
                     val, size=16,
                     bold=(r_idx == 0), color=color, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(6.0), Inches(12), Inches(0.6),
             "17× compression for ~27 dB PSNR; 3.4× compression for visually lossless 42 dB.",
             size=16, italic=True, color=ACCENT)

    # 14 - Rate-distortion
    s = blank_slide(prs)
    slide_header(s, "Rate–distortion trade-off", kicker="Results 3 of 3")
    s.shapes.add_picture(str(OUT / "fig_rate_distortion.png"),
                         Inches(0.7), Inches(1.3),
                         width=Inches(6.0))
    s.shapes.add_picture(str(OUT / "fig_bpp_vs_psnr.png"),
                         Inches(6.9), Inches(1.3),
                         width=Inches(6.0))
    add_text(s, Inches(0.7), Inches(5.9), Inches(12), Inches(0.6),
             "Concave PSNR vs. bpp curve — diminishing returns past ~1.2 bpp.",
             size=16, italic=True, color=MUTED)

    # 15 - Where is the error?
    s = blank_slide(prs)
    slide_header(s, "Where does the error live?", kicker="Diagnosis")
    s.shapes.add_picture(str(OUT / "fig_error_map.png"),
                         Inches(0.3), Inches(1.2),
                         width=Inches(12.7))
    add_text(s, Inches(0.7), Inches(6.2), Inches(12), Inches(0.6),
             "Reconstruction error is concentrated on edges — exactly where high-frequency "
             "DCT coefficients live, and exactly where the quantiser is coarsest.",
             size=14, italic=True, color=ACCENT)

    # 16 - Coefficient histogram
    s = blank_slide(prs)
    slide_header(s, "Why compression works — the statistics", kicker="Theory meets data")
    s.shapes.add_picture(str(OUT / "fig_coeff_histograms.png"),
                         Inches(0.5), Inches(1.3),
                         width=Inches(12.3))
    add_text(s, Inches(0.7), Inches(5.9), Inches(12), Inches(0.9),
             "86 % of quantised coefficients at Q=50 are exactly zero — a single Huffman "
             "symbol covers the bulk of the payload.",
             size=16, italic=True, color=MUTED)

    # 17 - Discussion / limitations
    s = blank_slide(prs)
    slide_header(s, "Limitations & honest caveats", kicker="Discussion")
    add_bullets(s, Inches(0.7), Inches(1.3), Inches(12), Inches(5.5), [
        "Compression ratios are based on the Shannon bound, not real Huffman coding.",
        "DC differentials and run-length markers are not modelled — small but real overhead.",
        "PSNR is only a loose proxy for perceived quality; SSIM would be more faithful.",
        "Blocking artefacts are a direct consequence of fixed 8×8 tiling — modern codecs use variable-size transforms plus in-loop deblocking.",
        "Only grayscale; extending to RGB would mean YCbCr conversion and a separate chroma table.",
    ])

    # 18 - Conclusion
    s = blank_slide(prs)
    slide_header(s, "What this project proves", kicker="Conclusion")
    add_bullets(s, Inches(0.7), Inches(1.3), Inches(12), Inches(5.5), [
        "Built a JPEG-baseline codec from scratch in ~250 lines of Python.",
        "Reproduced the textbook rate–distortion curve on a standard test image.",
        "Visualised every intermediate stage — DCT basis, quantised coefficients, error maps.",
        "Directly exercised Units 4 and 5 of IC2306, and Unit 2 via the DFT–DCT relationship.",
        "Natural next steps: add colour, add real entropy coding, measure SSIM.",
    ])

    # 19 - References / Q&A
    s = blank_slide(prs)
    add_rect(s, 0, 0, prs_width, prs_height, NAVY)
    add_text(s, Inches(0.6), Inches(2.0), Inches(12), Inches(1.2),
             "Thank you",
             size=54, bold=True, color=WHITE)
    add_text(s, Inches(0.6), Inches(3.2), Inches(12), Inches(0.8),
             "Questions?", size=28, italic=True, color=LIGHT)
    add_text(s, Inches(0.6), Inches(5.5), Inches(12), Inches(0.5),
             "References: Gonzalez & Woods (Ch. 7–8) • Proakis & Manolakis (Ch. 7) • "
             "Wallace, CACM 34(4), 1991 • Ahmed, Natarajan & Rao, IEEE TC, 1974",
             size=12, color=LIGHT)

    PPT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(PPT)
    print(f"wrote {PPT} — {len(prs.slides)} slides")


if __name__ == "__main__":
    build()
